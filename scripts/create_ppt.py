from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "LogicHR: 데이터로 증명하는 인사이트"
    subtitle.text = "\"감(Gut Feeling)이 아닌, SQL 로직으로 인사의 효율성을 증명하다\"\n\nTeam. 2215836-netizen"

    # Slide 1: Problem
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "1. 문제 정의 (Problem)"
    content = slide.placeholders[1]
    content.text = "현황: 많은 기업이 단순히 '매출'이 높거나 '야근'이 많은 부서를 고성과로 착각함.\n\nPain Point:\n"
    p = content.text_frame.add_paragraph()
    p.text = "• 인건비(투입) 대비 성과(산출)를 정량적으로 비교하기 어려움"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "• 특정 부서에서 불필요한 연장 근무(Leakage)가 발생해도 감지 못함"
    p.level = 1

    # Slide 2: Solution
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "2. 솔루션 (Solution)"
    content = slide.placeholders[1]
    content.text = "LogicHR: 데이터 기반 HR 의사결정 시스템\n"
    p = content.text_frame.add_paragraph()
    p.text = "• Core Engine: Python + SQLite (In-Memory DB)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "• Key Metric: 효율 지수 (Efficiency Index) & 인건비 누수 탐지"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "• Value: 복잡한 데이터 속에서 \"비용 효율성\"과 \"리스크\"를 즉시 추출"
    p.level = 1

    # Slide 3: Feature 1 - Efficiency Index
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "핵심 기능 1: 효율 지수 (Efficiency Index)"
    content = slide.placeholders[1]
    content.text = "\"단순한 성과 1등이 아니라, 가성비 1등을 찾습니다.\"\n"
    p = content.text_frame.add_paragraph()
    p.text = "• Logic: (목표 달성률 * 100) / (총 인건비 / 100만)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "• SQL 활용: JOIN으로 성과+근태 결합, Aggregation으로 비용 집계"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "• 효과: 적은 인원으로 높은 성과를 낸 '숨은 영웅' 부서 발굴"
    p.level = 1

    # Slide 4: Feature 2 - Leakage Detector
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "핵심 기능 2: 인건비 누수 탐지"
    content = slide.placeholders[1]
    content.text = "\"데이터 패턴으로 돈이 새는 구멍을 막습니다.\"\n"
    p = content.text_frame.add_paragraph()
    p.text = "• Logic: GROUP BY day_of_week & HAVING avg_hours > 9"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "• 인사이트: 특정 요일에만 야근이 폭증하지만 성과는 낮은 부서 경고(Alert)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "• 가치: 불필요한 초과 근무 수당 절감 및 업무 프로세스 개선"
    p.level = 1

    # Slide 5: Tech Stack
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "4. 기술 아키텍처 (Tech Stack)"
    content = slide.placeholders[1]
    content.text = "• Language: Python 3.9+"
    p = content.text_frame.add_paragraph()
    p.text = "• Limit: 외부 DB 서버 없이 로컬 파일/메모리에서 SQLite로 고성능 처리"
    p = content.text_frame.add_paragraph()
    p.text = "• Visualization: Streamlit, Plotly Express (Interactive Dashboard)"
    p = content.text_frame.add_paragraph()
    p.text = "• Collaboration: Git, GitHub"

    # Slide 6: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "5. 결론 및 기대 효과"
    content = slide.placeholders[1]
    content.text = "LogicHR 도입 시 기대 효과\n"
    p = content.text_frame.add_paragraph()
    p.text = "1. 비용 절감: 비효율적인 연장 근무 패턴 식별 및 개선"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "2. 공정한 평가: 투입 비용까지 고려한 입체적인 성과 평가 가능"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "3. 데이터 리터러시: 직관적인 대시보드와 SQL 공개로 조직 역량 강화"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "\n\"LogicHR은 단순한 대시보드가 아닙니다. 조직의 건강한 성장을 돕는 나침반입니다.\""
    p.level = 0
    p.font.bold = True

    prs.save('docs/LogicHR_Presentation.pptx')
    print("Presentation created successfully at docs/LogicHR_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
